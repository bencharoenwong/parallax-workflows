"""OOXML theme XML parsing + PPTX/DOCX extraction.

PPTX and DOCX share the same theme XML structure (drawingml namespace), so
the parser is shared. Theme XML provides canonical brand declarations
(colors, fonts) at high confidence (0.9). Body text is aggregated for the
voice corpus; voice extraction itself is LLM-driven in SKILL.md.
"""

from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List

from .voice import _voice_corpus_from_text


_OOXML_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _parse_ooxml_theme(theme_xml_bytes: bytes) -> Dict[str, Any]:
    """Parse OOXML theme XML and return color scheme + font scheme.

    Returns colors keyed by OOXML slot name (dk1, lt1, dk2, lt2, accent1..accent6,
    hlink, folHlink) and fonts keyed by role (major, minor).
    """
    from xml.etree import ElementTree as ET

    ns = {"a": _OOXML_A_NS}
    root = ET.fromstring(theme_xml_bytes)

    colors: Dict[str, str] = {}
    clr_scheme = root.find(".//a:clrScheme", ns)
    if clr_scheme is not None:
        for slot in clr_scheme:
            slot_name = slot.tag.split("}", 1)[-1]
            srgb = slot.find("a:srgbClr", ns)
            sys = slot.find("a:sysClr", ns)
            if srgb is not None and srgb.get("val"):
                colors[slot_name] = "#" + srgb.get("val").upper()
            elif sys is not None and sys.get("lastClr"):
                colors[slot_name] = "#" + sys.get("lastClr").upper()

    fonts: Dict[str, str] = {}
    font_scheme = root.find(".//a:fontScheme", ns)
    if font_scheme is not None:
        for role_elem_name, role_key in [("majorFont", "major"), ("minorFont", "minor")]:
            role_elem = font_scheme.find(f"a:{role_elem_name}", ns)
            if role_elem is not None:
                latin = role_elem.find("a:latin", ns)
                if latin is not None and latin.get("typeface"):
                    fonts[role_key] = latin.get("typeface")

    return {"colors": colors, "fonts": fonts}


def _theme_to_role_map(theme: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    """Map OOXML theme slots to white-label color roles.

    accent1 → primary, accent2 → secondary, accent3 → accent.
    dk1 → text, lt1 → background. Confidence 0.9 (theme XML is canonical).
    """
    raw = theme.get("colors", {})
    role_colors: Dict[str, Dict[str, Any]] = {}

    slot_to_role = [
        ("accent1", "primary"),
        ("accent2", "secondary"),
        ("accent3", "accent"),
        ("dk1", "text"),
        ("lt1", "background"),
    ]
    for slot, role in slot_to_role:
        if slot in raw:
            role_colors[role] = {"hex": raw[slot], "confidence": 0.9}

    fonts: Dict[str, Dict[str, Any]] = {}
    raw_fonts = theme.get("fonts", {})
    if "major" in raw_fonts:
        fonts["header"] = {"name": raw_fonts["major"], "confidence": 0.9}
    if "minor" in raw_fonts:
        fonts["body"] = {"name": raw_fonts["minor"], "confidence": 0.9}

    return {"colors": role_colors, "fonts": fonts}


def extract_from_pptx(pptx_path: str) -> Dict[str, Any]:
    """Extract branding (visual + voice corpus) from a PowerPoint file.

    Visual: theme XML provides precise colors/fonts (high confidence).
    Voice corpus: aggregated slide body text for LLM voice extraction in SKILL.md.
    """
    pptx_file = Path(pptx_path)
    base_return = {
        "colors": {},
        "logos": {},
        "fonts": {},
        "source": {"type": "pptx", "reference": pptx_path},
        "extracted_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "confidence_scores": {},
        "voice_corpus": {"text": "", "word_count": 0, "truncated": False},
    }

    if not pptx_file.exists():
        return {**base_return, "error": "PPTX file not found"}

    try:
        import zipfile
        theme: Dict[str, Any] = {"colors": {}, "fonts": {}}
        with zipfile.ZipFile(pptx_path, "r") as zf:
            theme_paths = [n for n in zf.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")]
            if theme_paths:
                with zf.open(theme_paths[0]) as f:
                    theme = _parse_ooxml_theme(f.read())

        role_map = _theme_to_role_map(theme)

        body_text = ""
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            chunks: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(r.text for r in para.runs).strip()
                            if line:
                                chunks.append(line)
            body_text = "\n".join(chunks)
        except Exception:
            pass

        corpus = _voice_corpus_from_text(body_text)

        confidence_scores: Dict[str, float] = {}
        for role, data in role_map["colors"].items():
            confidence_scores[f"color_{role}"] = data["confidence"]
        for usage, data in role_map["fonts"].items():
            confidence_scores[f"font_{usage}"] = data["confidence"]

        return {
            **base_return,
            "colors": role_map["colors"],
            "fonts": role_map["fonts"],
            "confidence_scores": confidence_scores,
            "voice_corpus": corpus,
        }

    except Exception as e:
        return {**base_return, "error": str(e)}


def extract_from_docx(docx_path: str) -> Dict[str, Any]:
    """Extract branding (visual + voice corpus) from a Word document.

    Visual: theme XML if present (high confidence), styles fallback otherwise.
    Voice corpus: aggregated paragraph text for LLM voice extraction in SKILL.md.
    """
    docx_file = Path(docx_path)
    base_return = {
        "colors": {},
        "logos": {},
        "fonts": {},
        "source": {"type": "docx", "reference": docx_path},
        "extracted_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "confidence_scores": {},
        "voice_corpus": {"text": "", "word_count": 0, "truncated": False},
    }

    if not docx_file.exists():
        return {**base_return, "error": "DOCX file not found"}

    try:
        import zipfile
        theme: Dict[str, Any] = {"colors": {}, "fonts": {}}
        with zipfile.ZipFile(docx_path, "r") as zf:
            theme_paths = [n for n in zf.namelist() if n.startswith("word/theme/") and n.endswith(".xml")]
            if theme_paths:
                with zf.open(theme_paths[0]) as f:
                    theme = _parse_ooxml_theme(f.read())

        role_map = _theme_to_role_map(theme)

        body_text = ""
        try:
            from docx import Document
            doc = Document(docx_path)
            chunks: List[str] = []
            for para in doc.paragraphs:
                line = para.text.strip()
                if line:
                    chunks.append(line)
            body_text = "\n".join(chunks)
        except Exception:
            pass

        corpus = _voice_corpus_from_text(body_text)

        confidence_scores: Dict[str, float] = {}
        for role, data in role_map["colors"].items():
            confidence_scores[f"color_{role}"] = data["confidence"]
        for usage, data in role_map["fonts"].items():
            confidence_scores[f"font_{usage}"] = data["confidence"]

        return {
            **base_return,
            "colors": role_map["colors"],
            "fonts": role_map["fonts"],
            "confidence_scores": confidence_scores,
            "voice_corpus": corpus,
        }

    except Exception as e:
        return {**base_return, "error": str(e)}
