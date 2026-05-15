"""Tests for PPTX/DOCX extractors, cross-source validation, merging,
and voice validation. Companion to test_extract.py / test_validator.py."""

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).parent.parent))

from extract import (
    cross_validate_visual,
    extract_from_docx,
    extract_from_pptx,
    merge_drafts,
    _assign_color_roles_by_frequency,
    _is_pure_white,
    _is_dark_text_candidate,
    _is_neutral_grey,
    _normalize_hex,
)
from validator import VoiceValidator


@pytest.fixture
def sample_pptx(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Test Asset Manager"
    body_shape = slide.shapes.add_textbox(0, 1000000, 5000000, 2000000)
    tf = body_shape.text_frame
    tf.text = "We invest with discipline. We avoid speculation."
    tf.add_paragraph().text = "Quarterly review covers global equities and macro positioning."

    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def sample_docx(tmp_path):
    from docx import Document

    doc = Document()
    doc.add_heading("Quarterly Letter", 1)
    doc.add_paragraph("Dear clients, our portfolio enters Q2 with measured optimism.")
    doc.add_paragraph("We continue to favor quality over momentum.")

    path = tmp_path / "sample.docx"
    doc.save(str(path))
    return str(path)


class TestExtractFromPptx:
    def test_returns_expected_keys(self, sample_pptx):
        draft = extract_from_pptx(sample_pptx)
        for key in ["colors", "logos", "fonts", "source", "extracted_at",
                    "confidence_scores", "voice_corpus"]:
            assert key in draft, f"missing key: {key}"

    def test_source_type_is_pptx(self, sample_pptx):
        draft = extract_from_pptx(sample_pptx)
        assert draft["source"]["type"] == "pptx"
        assert draft["source"]["reference"] == sample_pptx

    def test_theme_colors_extracted(self, sample_pptx):
        draft = extract_from_pptx(sample_pptx)
        assert "primary" in draft["colors"], "expected accent1→primary mapping"
        assert draft["colors"]["primary"]["hex"].startswith("#")
        assert draft["colors"]["primary"]["confidence"] == 0.9

    def test_theme_fonts_extracted(self, sample_pptx):
        draft = extract_from_pptx(sample_pptx)
        assert "header" in draft["fonts"]
        assert "body" in draft["fonts"]
        assert draft["fonts"]["header"]["confidence"] == 0.9

    def test_voice_corpus_populated(self, sample_pptx):
        draft = extract_from_pptx(sample_pptx)
        assert draft["voice_corpus"]["word_count"] > 0
        assert "discipline" in draft["voice_corpus"]["text"].lower()

    def test_missing_file_returns_error(self, tmp_path):
        draft = extract_from_pptx(str(tmp_path / "nonexistent.pptx"))
        assert "error" in draft
        assert draft["voice_corpus"]["word_count"] == 0
        assert draft["colors"] == {}


class TestExtractFromDocx:
    def test_returns_expected_keys(self, sample_docx):
        draft = extract_from_docx(sample_docx)
        for key in ["colors", "logos", "fonts", "source", "extracted_at",
                    "confidence_scores", "voice_corpus"]:
            assert key in draft

    def test_source_type_is_docx(self, sample_docx):
        draft = extract_from_docx(sample_docx)
        assert draft["source"]["type"] == "docx"

    def test_theme_extracted(self, sample_docx):
        draft = extract_from_docx(sample_docx)
        assert len(draft["colors"]) > 0
        assert len(draft["fonts"]) > 0

    def test_voice_corpus_includes_paragraphs(self, sample_docx):
        draft = extract_from_docx(sample_docx)
        text = draft["voice_corpus"]["text"].lower()
        assert "quarterly" in text
        assert "momentum" in text

    def test_missing_file_returns_error(self, tmp_path):
        draft = extract_from_docx(str(tmp_path / "nonexistent.docx"))
        assert "error" in draft


class TestCrossValidateVisual:
    def test_single_draft_returns_empty(self, sample_pptx):
        draft = extract_from_pptx(sample_pptx)
        result = cross_validate_visual([draft])
        assert result["mismatches"] == []
        assert result["agreements"] == []

    def test_two_identical_drafts_all_agreements(self, sample_pptx):
        d1 = extract_from_pptx(sample_pptx)
        d2 = extract_from_pptx(sample_pptx)
        result = cross_validate_visual([d1, d2])
        assert result["mismatches"] == []
        assert len(result["agreements"]) > 0

    def test_color_mismatch_flagged(self):
        d1 = {
            "colors": {"primary": {"hex": "#FF0000", "confidence": 0.9}},
            "fonts": {},
            "source": {"reference": "a"},
        }
        d2 = {
            "colors": {"primary": {"hex": "#00FF00", "confidence": 0.9}},
            "fonts": {},
            "source": {"reference": "b"},
        }
        result = cross_validate_visual([d1, d2])
        assert len(result["mismatches"]) == 1
        assert result["mismatches"][0]["field"] == "colors.primary"
        sources = {v["source"] for v in result["mismatches"][0]["values"]}
        assert sources == {"a", "b"}

    def test_font_mismatch_case_insensitive(self):
        d1 = {"colors": {}, "fonts": {"header": {"name": "Calibri", "confidence": 0.9}}, "source": {"reference": "a"}}
        d2 = {"colors": {}, "fonts": {"header": {"name": "calibri", "confidence": 0.9}}, "source": {"reference": "b"}}
        result = cross_validate_visual([d1, d2])
        assert result["mismatches"] == []
        assert any(a["field"] == "fonts.header" for a in result["agreements"])


class TestMergeDrafts:
    def test_empty_list_returns_empty_dict(self):
        assert merge_drafts([]) == {}

    def test_single_draft_returns_unchanged(self, sample_pptx):
        d = extract_from_pptx(sample_pptx)
        assert merge_drafts([d]) is d

    def test_higher_confidence_wins(self):
        d1 = {
            "colors": {"primary": {"hex": "#AAAAAA", "confidence": 0.5}},
            "logos": {}, "fonts": {}, "confidence_scores": {},
            "source": {"reference": "a"}, "voice_corpus": {"text": "", "word_count": 0},
        }
        d2 = {
            "colors": {"primary": {"hex": "#BBBBBB", "confidence": 0.9}},
            "logos": {}, "fonts": {}, "confidence_scores": {},
            "source": {"reference": "b"}, "voice_corpus": {"text": "", "word_count": 0},
        }
        merged = merge_drafts([d1, d2])
        assert merged["colors"]["primary"]["hex"] == "#BBBBBB"

    def test_voice_corpora_concatenated(self, sample_pptx, sample_docx):
        d1 = extract_from_pptx(sample_pptx)
        d2 = extract_from_docx(sample_docx)
        merged = merge_drafts([d1, d2])
        assert merged["voice_corpus"]["word_count"] == (
            d1["voice_corpus"]["word_count"] + d2["voice_corpus"]["word_count"]
        )
        assert "by_source" in merged["voice_corpus"]
        assert len(merged["voice_corpus"]["by_source"]) == 2
        assert merged["source"]["type"] == "multi"


class TestVoiceValidator:
    def test_disabled_returns_skipped(self):
        result = VoiceValidator.validate_voice({"enabled": False})
        assert result["status"] == "skipped"

    def test_corpus_below_min_fails(self):
        voice = {
            "enabled": True,
            "positioning": "X.",
            "tone": {"register": "r", "primary_attributes": ["a", "b", "c"]},
            "core_rules": ["x", "y"],
            "anti_filler": ["a", "b", "c"],
            "source_corpus": {"word_count": 100},
        }
        result = VoiceValidator.validate_voice(voice)
        assert result["status"] == "fail"
        assert result["checks"]["corpus"]["status"] == "fail"

    def test_corpus_below_recommended_warns(self):
        voice = {
            "enabled": True,
            "positioning": "X.",
            "tone": {"register": "r", "primary_attributes": ["a", "b", "c"]},
            "core_rules": ["x", "y"],
            "anti_filler": ["a", "b", "c"],
            "source_corpus": {"word_count": 1000},
        }
        result = VoiceValidator.validate_voice(voice)
        assert result["status"] == "warn"
        assert result["checks"]["corpus"]["status"] == "warn"

    def test_full_pass(self):
        voice = {
            "enabled": True,
            "positioning": "Disciplined institutional asset manager.",
            "tone": {
                "register": "formal-institutional",
                "primary_attributes": ["measured", "evidence-led", "client-first"],
                "avoid_attributes": ["hyperbolic"],
            },
            "core_rules": ["Never speculate", "Always cite evidence"],
            "anti_filler": ["leverage", "synergy", "best-in-class"],
            "source_corpus": {"word_count": 2500},
        }
        result = VoiceValidator.validate_voice(voice)
        assert result["status"] == "pass"

class TestFolderModeIntegration:
    """Integration test for SKILL.md folder mode (F-1 to F-4 flow).

    The semantic classification step (F-2) is LLM-driven; this test exercises
    the F-4 mechanical extraction pattern: iterate files, dispatch to the
    right extractor by extension, merge drafts, cross-validate, append
    text-only content to the voice corpus.
    """

    def test_mixed_folder_extracts_and_merges(self, tmp_path, sample_pptx, sample_docx):
        from pptx import Presentation
        from docx import Document

        # Build a folder with PPTX, DOCX, and a plain .txt (voice-only contributor)
        folder = tmp_path / "client_collateral"
        folder.mkdir()

        # Re-use the existing fixtures via shutil-like copy (write fresh files)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Quarterly Newsletter Q1"
        body = slide.shapes.add_textbox(0, 1000000, 5000000, 2000000).text_frame
        body.text = "Our discipline serves clients across credit and equity. We avoid speculation."
        pptx_path = folder / "newsletter-q1.pptx"
        prs.save(str(pptx_path))

        doc = Document()
        doc.add_heading("Client Letter", 1)
        doc.add_paragraph("Dear clients, we maintain measured optimism amid uncertainty.")
        doc.add_paragraph("Quality remains our anchor through cycles.")
        docx_path = folder / "client-letter-jan.docx"
        doc.save(str(docx_path))

        # Voice-only text file (extension .md / .txt are read directly per F-4 flow)
        txt_path = folder / "research-notes.md"
        txt_path.write_text(
            "# Internal Research Note\n\n"
            "We continue to favor capital-intensive businesses with hard-asset moats. "
            "The structural shift toward physical scarcity is well underway.",
            encoding="utf-8",
        )

        # F-4 mechanical pattern (LLM does F-1 to F-3 semantic classification first;
        # this test pins the deterministic extraction phase):
        drafts = [
            extract_from_pptx(str(pptx_path)),
            extract_from_docx(str(docx_path)),
        ]

        # Both OOXML drafts came through cleanly
        assert all(d.get("error") is None for d in drafts), \
            f"unexpected error in drafts: {[d.get('error') for d in drafts]}"
        assert all(d["voice_corpus"]["word_count"] > 0 for d in drafts)

        # Merge OOXML drafts
        merged = merge_drafts(drafts)
        assert merged["source"]["type"] == "multi"
        assert "by_source" in merged["voice_corpus"]
        assert len(merged["voice_corpus"]["by_source"]) == 2

        # Cross-validate (default Office theme — should agree on visual fields)
        xv = cross_validate_visual(drafts)
        # Expect agreements on at least colors.primary (both default to Office accent1)
        agreement_fields = {a["field"] for a in xv["agreements"]}
        assert "colors.primary" in agreement_fields, \
            f"expected colors.primary agreement, got {agreement_fields}"

        # Append the text-only file's content to voice_corpus (per SKILL.md F-4)
        extra_text = txt_path.read_text(encoding="utf-8")
        before_words = merged["voice_corpus"]["word_count"]
        merged["voice_corpus"]["text"] += "\n\n" + extra_text
        merged["voice_corpus"]["word_count"] = before_words + len(extra_text.split())

        assert merged["voice_corpus"]["word_count"] > before_words
        assert "physical scarcity" in merged["voice_corpus"]["text"].lower(), \
            "text-only file content not propagated into voice corpus"

    def test_folder_with_single_file_no_merge_needed(self, tmp_path, sample_pptx):
        """Single-file folder: extract returns the draft as-is; no merge_drafts call."""
        # Just one supported file — merge_drafts([draft]) returns the draft unchanged
        draft = extract_from_pptx(sample_pptx)
        merged = merge_drafts([draft])
        assert merged is draft  # identity preservation
        assert merged["source"]["type"] == "pptx"  # NOT promoted to "multi"


class TestColorClassifiers:
    def test_normalize_hex_3_to_6(self):
        assert _normalize_hex("#fff") == "#FFFFFF"
        assert _normalize_hex("#abc") == "#AABBCC"
        assert _normalize_hex("#123456") == "#123456"
        assert _normalize_hex("ff0000") == "#FF0000"

    def test_pure_white(self):
        assert _is_pure_white("#FFFFFF")
        assert _is_pure_white("#fff")
        assert not _is_pure_white("#FEFEFE")
        assert not _is_pure_white("#000000")

    def test_dark_text_candidate(self):
        assert _is_dark_text_candidate("#000000")
        assert _is_dark_text_candidate("#1A1A1A")
        assert _is_dark_text_candidate("#333333")  # max=51, < 80
        assert not _is_dark_text_candidate("#5A597A")  # max=122
        assert not _is_dark_text_candidate("#FFFFFF")

    def test_neutral_grey(self):
        assert _is_neutral_grey("#808080")     # mid-grey
        assert _is_neutral_grey("#CCCCCC")     # light grey
        assert _is_neutral_grey("#9A9A9A")
        assert not _is_neutral_grey("#FFFFFF") # pure white excluded
        assert not _is_neutral_grey("#000000") # pure black excluded
        assert not _is_neutral_grey("#5A597A") # saturated brand color
        assert not _is_neutral_grey("#FF0000") # saturated


class TestAssignColorRolesByFrequency:
    def test_pure_white_becomes_background(self):
        colors = [{"hex": "#FFFFFF", "confidence": 0.95}] * 10
        result = _assign_color_roles_by_frequency(colors)
        assert result.get("background", {}).get("hex") == "#FFFFFF"
        assert "primary" not in result  # white must NOT be assigned to primary

    def test_dark_color_becomes_text(self):
        colors = [{"hex": "#FFFFFF", "confidence": 0.95}] * 5 + [{"hex": "#1A1A1A", "confidence": 0.95}] * 3
        result = _assign_color_roles_by_frequency(colors)
        assert result.get("text", {}).get("hex") == "#1A1A1A"

    def test_most_frequent_brand_color_becomes_primary(self):
        colors = (
            [{"hex": "#FFFFFF", "confidence": 0.95}] * 50  # background, ignored for brand
            + [{"hex": "#5A597A", "confidence": 0.95}] * 8  # primary candidate
            + [{"hex": "#676C85", "confidence": 0.95}] * 5  # secondary candidate
            + [{"hex": "#A0A0A0", "confidence": 0.95}] * 4  # neutral grey, excluded
        )
        result = _assign_color_roles_by_frequency(colors)
        assert result["primary"]["hex"] == "#5A597A"
        assert result["secondary"]["hex"] == "#676C85"
        # neutral grey should not appear as accent
        assert result.get("accent", {}).get("hex") != "#A0A0A0"

    def test_confidence_scales_with_frequency(self):
        many = [{"hex": "#FF0000", "confidence": 0.95}] * 10
        few = [{"hex": "#00FF00", "confidence": 0.95}] * 1
        result_many = _assign_color_roles_by_frequency(many)
        result_few = _assign_color_roles_by_frequency(few)
        assert result_many["primary"]["confidence"] == 0.85   # ≥5
        assert result_few["primary"]["confidence"] == 0.6     # singleton

    def test_off_white_backgrounds_not_routed_to_primary(self):
        """Regression for review finding: #F5F5F5, #FAFAFA, #F0F0F0 should be
        treated as neutral greys (excluded from brand pool), not assigned to
        primary. Brand color in saturated palette should win."""
        # Common website pattern: bulk of CSS uses off-white shades + a few
        # saturated brand colors.
        colors = (
            [{"hex": "#F5F5F5", "confidence": 0.95}] * 30   # background-ish, off-white
            + [{"hex": "#FAFAFA", "confidence": 0.95}] * 20  # background-ish
            + [{"hex": "#F0F0F0", "confidence": 0.95}] * 15  # background-ish
            + [{"hex": "#5A597A", "confidence": 0.95}] * 5   # actual brand color (singleton-ish)
        )
        result = _assign_color_roles_by_frequency(colors)
        # The saturated brand color should win primary, NOT the most-frequent off-white
        assert result.get("primary", {}).get("hex") == "#5A597A"
        # Off-white shades should not appear as any brand role
        for role in ("primary", "secondary", "accent"):
            assigned = result.get(role, {}).get("hex")
            if assigned is not None:
                assert assigned not in ("#F5F5F5", "#FAFAFA", "#F0F0F0"), \
                    f"{role} should not be off-white {assigned}"

    def test_three_digit_hex_normalized(self):
        colors = [{"hex": "#fff", "confidence": 0.95}] * 5 + [{"hex": "#f00", "confidence": 0.95}] * 5
        result = _assign_color_roles_by_frequency(colors)
        assert result["background"]["hex"] == "#FFFFFF"
        assert result["primary"]["hex"] == "#FF0000"

    def test_empty_input_returns_empty(self):
        assert _assign_color_roles_by_frequency([]) == {}


class TestVoiceValidatorMissingFields:
    def test_missing_fields_warn(self):
        voice = {
            "enabled": True,
            "positioning": "",  # missing
            "tone": {"register": "", "primary_attributes": []},  # missing
            "core_rules": ["x", "y"],
            "anti_filler": ["a", "b", "c"],
            "source_corpus": {"word_count": 2500},
        }
        result = VoiceValidator.validate_voice(voice)
        assert result["status"] == "fail"
        assert "positioning" in result["checks"]["completeness"]["missing"]


def test_ooxml_pptx_typography_and_radii(tmp_path):
    import zipfile
    from extract.ooxml import extract_from_pptx

    pptx_path = tmp_path / "test.pptx"
    
    master_xml = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
    <p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        <p:txStyles>
            <p:titleStyle><a:lvl1pPr><a:defRPr sz="4400" b="1"/></a:lvl1pPr></p:titleStyle>
            <p:bodyStyle><a:lvl5pPr><a:defRPr sz="1200"/><a:lnSpc><a:spcPct val="150000"/></a:lnSpc></a:lvl5pPr></p:bodyStyle>
        </p:txStyles>
    </p:sldMaster>'''
    
    slide_xml = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
    <p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        <p:spTree>
            <p:sp><p:spPr><a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 4000"/></a:avLst></a:prstGeom></p:spPr></p:sp>
            <p:sp><p:spPr><a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 4000"/></a:avLst></a:prstGeom></p:spPr></p:sp>
            <p:sp><p:spPr><a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 12000"/></a:avLst></a:prstGeom></p:spPr></p:sp>
        </p:spTree>
    </p:sld>'''

    with zipfile.ZipFile(pptx_path, "w") as zf:
        zf.writestr("ppt/slideMasters/slideMaster1.xml", master_xml)
        zf.writestr("ppt/slides/slide1.xml", slide_xml)
        
    draft = extract_from_pptx(str(pptx_path))
    
    assert "typography" in draft
    # DESIGN.md spec only accepts px / rem / em — extractor converts pt → px
    # at 96dpi (1pt = 4/3 px). 44pt → 59px, 12pt → 16px. letterSpacing must
    # carry a unit ("0em"), not bare "0".
    assert draft["typography"]["h1"]["fontSize"] == "59px"
    assert draft["typography"]["h1"]["fontWeight"] == 700
    assert draft["typography"]["h1"]["letterSpacing"] == "0em"
    assert draft["typography"]["body-md"]["fontSize"] == "16px"
    assert draft["typography"]["body-md"]["lineHeight"] == "1.5"
    assert "h2" not in draft["typography"]
    
    assert "rounded" in draft
    assert draft["rounded"]["sm"] == "4px"
    assert draft["rounded"]["md"] == "12px"
    
    assert draft["confidence_scores"]["typography.h1"] == 0.85
    assert draft["confidence_scores"]["typography.body-md"] == 0.85
    assert draft["confidence_scores"]["rounded"] == 0.5

def test_ooxml_docx_typography(tmp_path):
    import zipfile
    from extract.ooxml import extract_from_docx

    docx_path = tmp_path / "test.docx"

    styles_xml = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
    <w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
        <w:style w:type="paragraph" w:styleId="Heading1">
            <w:name w:val="heading 1"/>
            <w:rPr><w:sz w:val="48"/><w:b/></w:rPr>
            <w:pPr><w:spacing w:line="360"/></w:pPr>
        </w:style>
    </w:styles>'''

    with zipfile.ZipFile(docx_path, "w") as zf:
        zf.writestr("word/styles.xml", styles_xml)

    draft = extract_from_docx(str(docx_path))
    
    assert "typography" in draft
    # DOCX half-points → pt → px. w:sz=48 → 24pt → 32px.
    assert draft["typography"]["h1"]["fontSize"] == "32px"
    assert draft["typography"]["h1"]["fontWeight"] == 700
    assert draft["typography"]["h1"]["letterSpacing"] == "0em"
    assert draft["typography"]["h1"]["lineHeight"] == "1.5"
    assert draft["confidence_scores"]["typography.h1"] == 0.85

def test_ooxml_no_round_rect_yields_empty(tmp_path):
    from unittest.mock import patch, MagicMock
    import zipfile
    from extract.ooxml import extract_from_pptx

    pptx_path = tmp_path / "test.pptx"
    pptx_path.write_text("fake")

    mock_zip = MagicMock()
    mock_zip.namelist.return_value = ["ppt/slides/slide1.xml"]
    
    slide_xml = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
    <p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
    </p:sld>"""

    def mock_open(name):
        m = MagicMock()
        m = io.BytesIO(slide_xml)
        
        return m
    
    mock_zip.open = mock_open
    
    with patch("zipfile.ZipFile", return_value=mock_zip):
        draft = extract_from_pptx(str(pptx_path))
        assert "rounded" not in draft

