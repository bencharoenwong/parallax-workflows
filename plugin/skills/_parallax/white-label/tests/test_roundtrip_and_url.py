"""Save → reload roundtrip + urllib URL-fallback regression tests.

Closes the integration gaps identified by the council audit:
  - Test 1: end-to-end draft → build_config_from_draft → yaml dump → load_client_branding
  - Test 3: urllib fallback path for URL extraction (mocked HTTP, no live network)
"""

import importlib
import importlib.util
import sys
from pathlib import Path
from unittest import mock

import pytest
import yaml

# Test files use sys.path manipulation; conftest.py loads loader as a module
sys.path.insert(0, str(Path(__file__).parent.parent))

from extract import (  # noqa: E402
    extract_from_pptx,
    extract_from_url,
)
import extract.web_pdf as web_pdf_module  # noqa: E402


# Load loader.py via conftest's existing pattern
HERE = Path(__file__).parent
LOADER_PATH = HERE.parent / "loader.py"
spec = importlib.util.spec_from_file_location("loader", LOADER_PATH)
loader = importlib.util.module_from_spec(spec)
spec.loader.exec_module(loader)


@pytest.fixture
def sample_pptx(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Test Asset Manager"
    body_shape = slide.shapes.add_textbox(0, 1000000, 5000000, 2000000)
    tf = body_shape.text_frame
    tf.text = "We invest with discipline. We avoid speculation."
    tf.add_paragraph().text = "Quarterly review covers global equities."

    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return str(path)


# ---------------------------------------------------------------------------
# Test 1: Save → reload roundtrip
# ---------------------------------------------------------------------------


class TestSaveReloadRoundtrip:
    def test_pptx_draft_roundtrips_through_yaml(self, tmp_path, sample_pptx, monkeypatch):
        """Extract → build_config → yaml.dump → load_client_branding preserves data."""
        # 1. Extract
        draft = extract_from_pptx(sample_pptx)

        # 2. Build config
        config = loader.build_config_from_draft(
            draft,
            client_name="Test Asset Manager",
            extracted_by="test@example.com",
        )

        # Required structure present
        assert "metadata" in config
        assert "branding" in config
        assert "confidence_scores" in config
        assert config["voice"] == {"enabled": False}  # voice not extracted in this test

        # 3. Write to tempfile
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config, sort_keys=False))

        # 4. Load via loader (monkeypatch the path)
        monkeypatch.setattr(loader, "_CONFIG_PATH", config_path)
        monkeypatch.setattr(loader, "_SCHEMA", loader._JSONSCHEMA)
        result = loader.load_client_branding()

        # 5. Assert roundtrip integrity
        # Schema validation passed (no schema_invalid error)
        assert result["error"] is None or "logo_missing" in (result["error"] or ""), \
            f"Unexpected error: {result['error']}"

        # Colors round-tripped (PPTX default Office theme has accent1 = #4F81BD)
        assert result["colors"]["primary"] == "#4F81BD"
        assert result["colors"]["background"] == "#FFFFFF"

        # Source type preserved
        assert result["source"]["type"] == "pptx"

        # Voice section preserved (disabled, not extracted)
        assert result["voice"] == {"enabled": False}

        # Multi-source absent (single-source extraction)
        assert result["multi_source"] == {}

    def test_voice_enabled_roundtrips(self, tmp_path, sample_pptx, monkeypatch):
        """Draft with voice section round-trips through yaml without losing fields."""
        draft = extract_from_pptx(sample_pptx)

        # Inject a voice section as Step 1.5 would
        draft["voice"] = {
            "enabled": True,
            "positioning": "Disciplined institutional asset manager.",
            "tone": {
                "register": "formal-institutional",
                "primary_attributes": ["measured", "evidence-led", "client-first"],
                "avoid_attributes": ["hyperbolic"],
            },
            "core_rules": ["Never speculate", "Always cite evidence"],
            "anti_filler": ["leverage", "synergy", "best-in-class"],
            "audience_adaptation": [],
            "channel_notes": [],
            "drafted_vs_sent": [],
            "company_context": "We are a long-only credit and equity manager.",
            "disclaimers": [
                {"jurisdiction": "MAS", "text": "Regulated by MAS.", "placement": "footer"}
            ],
            "source_corpus": {
                "documents": [sample_pptx],
                "word_count": 2500,
                "confidence": 0.85,
                "notes": "",
            },
        }

        config = loader.build_config_from_draft(draft, client_name="Test Co")
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config, sort_keys=False))

        monkeypatch.setattr(loader, "_CONFIG_PATH", config_path)
        monkeypatch.setattr(loader, "_SCHEMA", loader._JSONSCHEMA)
        result = loader.load_client_branding()

        assert result["voice"]["enabled"] is True
        assert result["voice"]["positioning"] == "Disciplined institutional asset manager."
        assert result["voice"]["tone"]["register"] == "formal-institutional"
        assert "measured" in result["voice"]["tone"]["primary_attributes"]
        assert len(result["voice"]["core_rules"]) == 2
        assert len(result["voice"]["anti_filler"]) == 3
        assert result["voice"]["disclaimers"][0]["jurisdiction"] == "MAS"
        assert result["voice"]["source_corpus"]["word_count"] == 2500

    def test_multi_source_roundtrips(self, tmp_path, sample_pptx, monkeypatch):
        """Multi-source draft preserves mismatches/agreements through save+load."""
        draft = extract_from_pptx(sample_pptx)
        draft["multi_source"] = {
            "sources": [
                {"type": "url", "reference": "https://example.com"},
                {"type": "pptx", "reference": sample_pptx},
            ],
            "mismatches": [
                {
                    "field": "fonts.body",
                    "values": [
                        {"source": "https://example.com", "value": "calibri"},
                        {"source": sample_pptx, "value": "cambria"},
                    ],
                }
            ],
            "agreements": [{"field": "colors.primary", "value": "#4F81BD"}],
        }

        config = loader.build_config_from_draft(draft, client_name="Test")
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config, sort_keys=False))

        monkeypatch.setattr(loader, "_CONFIG_PATH", config_path)
        monkeypatch.setattr(loader, "_SCHEMA", loader._JSONSCHEMA)
        result = loader.load_client_branding()

        assert result["multi_source"]["sources"]
        assert len(result["multi_source"]["mismatches"]) == 1
        assert result["multi_source"]["mismatches"][0]["field"] == "fonts.body"
        assert len(result["multi_source"]["agreements"]) == 1


# ---------------------------------------------------------------------------
# Test 3: urllib fallback regression
# ---------------------------------------------------------------------------


# A trimmed but realistic asset-manager homepage HTML stub. Pins the expected
# extraction outputs. If the URL extraction logic changes shape, this fixture
# catches the regression.
_BRAND_STUB_HTML = b"""
<!DOCTYPE html>
<html>
<head>
<style>
body { background: #FFFFFF; color: #333333; font-family: 'Calibri', sans-serif; }
.brand { color: #5A597A; }
.brand-secondary { color: #5A597A; }
.brand-accent { color: #676C85; }
.heading { color: #5A597A; }
.section { border-color: #5A597A; }
</style>
</head>
<body>
<header>
<img src="https://www.example.com/assets/brand/logo.png" alt="Example Asset Management logo">
</header>
<main>
<h1>Lorem ipsum dolor sit amet, consectetur adipiscing elit.</h1>
<p>Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua.</p>
<p>Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris.</p>
<p>Duis aute irure dolor in reprehenderit in voluptate velit esse.</p>
</main>
</body>
</html>
"""


class _StubResponse:
    """Minimal stand-in for urllib's HTTPResponse.

    NOTE: bare `mock.MagicMock()` for `.headers` is a footgun. The production
    code in _fetch_external_stylesheets does:

        ctype = resp.headers.get_content_type() if hasattr(...) else ...
        if ctype and "css" not in ctype.lower() and "text" not in ctype.lower():
            continue

    With a bare MagicMock, hasattr is always True and `ctype` becomes a
    MagicMock instance whose `__contains__` returns False — so `"css" not in
    ctype` is True and the content-type guard silently skips the body. Use
    real lambdas via spec= so the guard exercises the path it claims to.
    """
    def __init__(self, body: bytes, content_type: str = "text/html", final_url: str = "https://example.com/"):
        self._body = body
        self._final_url = final_url
        self.read_calls = 0
        self.headers = mock.MagicMock(spec=["get_content_charset", "get_content_type", "get"])
        self.headers.get_content_charset = lambda: "utf-8"
        self.headers.get_content_type = lambda: content_type
        self.headers.get = lambda key, default=None: content_type if key.lower() == "content-type" else default

    def read(self, amt=None):
        self.read_calls += 1
        # Real urllib HTTPResponse.read accepts an optional size cap.
        if amt is None:
            return self._body
        return self._body[:amt]

    def __enter__(self):
        return self

    def __exit__(self, *args):
        return False

    def geturl(self):
        return self._final_url

    def close(self):
        pass


class TestUrlFallbackRegression:
    @pytest.fixture(autouse=True)
    def public_dns_only(self, monkeypatch):
        """All URL tests are hermetic and resolve fixture hosts to a public IP."""
        def fake_getaddrinfo(host, port, *args, **kwargs):
            return [(2, 1, 6, "", ("93.184.216.34", port))]

        monkeypatch.setattr("socket.getaddrinfo", fake_getaddrinfo)

    def test_public_url_extracts_brand_signature(self, monkeypatch):
        """The validated URL opener's HTML supplies the brand signature."""
        # Intercept urllib at the import site used by extract.web_pdf.
        def fake_urlopen(req, timeout=None):
            return _StubResponse(_BRAND_STUB_HTML)

        with mock.patch("extract.web_pdf._network_open", fake_urlopen):
            draft = extract_from_url("https://www.example.com/")

        # Logo extracted (the brand logo, not a strategy image)
        assert "primary" in draft["logos"]
        assert "logo.png" in draft["logos"]["primary"]["url"]

        # Background detected via frequency-based role assignment
        assert draft["colors"].get("background", {}).get("hex") == "#FFFFFF"

        # Text detected as a dark color
        text_color = draft["colors"].get("text", {}).get("hex", "")
        assert text_color.startswith("#"), "text color should be assigned"

        # The dominant brand color (#5A597A appears 5+ times in the stub) is primary
        assert draft["colors"]["primary"]["hex"] == "#5A597A"

        # No error
        assert draft.get("error") is None

    def test_voice_corpus_excludes_script_and_style_bodies(self):
        """Tag-stripping alone keeps <script>/<style>/JSON-LD *contents*, which
        on a real page fill the corpus's leading-token window with minified JS
        while word_count still looks healthy."""
        page_html = b"""
<!DOCTYPE html><html><head>
<style>.brand{content:"noisetokenCss";font-family:Inter}</style>
<script>var noisetokenJs=function(){return 1};</script>
<script type="application/ld+json">{"@type":"Org","name":"noisetokenJsonLd"}</script>
<template><span>noisetokenTemplate</span></template>
<!-- noisetokenComment -->
</head><body><p>Disciplined capital allocation for institutional investors.</p>
</body></html>
"""

        def fake_urlopen(req, timeout=None):
            return _StubResponse(page_html, content_type="text/html")

        with mock.patch("extract.web_pdf._network_open", fake_urlopen):
            draft = extract_from_url("https://example.com/")

        corpus = draft["voice_corpus"]["text"]
        for noise in ("noisetokenCss", "noisetokenJs", "noisetokenJsonLd",
                      "noisetokenTemplate", "noisetokenComment"):
            assert noise not in corpus, f"{noise} leaked into the voice corpus"
        assert "Disciplined capital allocation" in corpus

    @pytest.mark.parametrize("tag", ["script", "style", "template", "noscript"])
    def test_non_prose_stripping_is_linear_on_unterminated_openers(self, tag):
        """A lazy DOTALL regex re-scans to end-of-document for every opener
        that never closes, so a hostile page of bare openers hangs the
        session. The scan must stay linear."""
        import time

        hostile = f"<{tag}>" * 40_000 + "x" * 400_000
        start = time.monotonic()
        out = web_pdf_module._strip_non_prose_blocks(hostile)
        assert time.monotonic() - start < 2.0
        assert out.strip() == ""

    def test_non_prose_stripping_preserves_lookalike_tag_names(self):
        """`<scripture>` shares a prefix with `<script>` but is not a
        non-prose block; only a real tag-name boundary may trigger a strip."""
        out = web_pdf_module._strip_non_prose_blocks(
            "<scripture>keepme</scripture><script>dropme</script>")
        assert "keepme" in out
        assert "dropme" not in out

    def test_non_prose_stripping_survives_unicode_case_expansion(self):
        """U+0130 lowercases to TWO code points, so offsets taken on a
        case-folded copy drift against the source and slice mid-character.
        Turkish brand prose ('İstanbul', 'İş Bankası') hits this on every page.
        """
        out = web_pdf_module._strip_non_prose_blocks(
            "İstanbul İş Bankası İzmir keep <script>DROPME</script> tail")
        assert "DROPME" not in out
        assert "<" not in out and ">" not in out
        assert "İstanbul İş Bankası İzmir keep" in out
        assert "tail" in out

    def test_voice_corpus_strips_scripts_on_unicode_heavy_page(self):
        page_html = (
            "<html><body><p>"
            + "İş " * 200
            + "</p><script>var noisetokenJs=1;</script>"
            "<p>İstanbul capital allocation.</p></body></html>"
        ).encode("utf-8")

        def fake_urlopen(req, timeout=None):
            return _StubResponse(page_html, content_type="text/html")

        with mock.patch("extract.web_pdf._network_open", fake_urlopen):
            draft = extract_from_url("https://example.com/")

        corpus = draft["voice_corpus"]["text"]
        assert "noisetokenJs" not in corpus
        assert "İstanbul capital allocation." in corpus

    def test_unterminated_opener_consumes_the_remaining_document(self):
        out = web_pdf_module._strip_non_prose_blocks(
            "<p>keepme</p><script>var x = '</p>not prose'")
        assert "keepme" in out
        assert "not prose" not in out

    def test_non_ssrf_failure_is_not_labelled_a_destination_rejection(
            self, monkeypatch):
        """Only a public-address rejection may report 'not public'. A plain
        ValueError from any downstream extractor must not send the operator
        debugging DNS."""
        monkeypatch.setattr(
            web_pdf_module, "_assign_color_roles_by_frequency",
            mock.Mock(side_effect=ValueError("unrelated parse failure")))

        def fake_urlopen(req, timeout=None):
            return _StubResponse(b"<html><body>hi</body></html>",
                                 content_type="text/html")

        with mock.patch("extract.web_pdf._network_open", fake_urlopen):
            draft = extract_from_url("https://example.com/")

        assert draft["error"] == "URL extraction failed"

    def test_external_stylesheet_following_recovers_fonts(self):
        """When the page links to external CSS via <link rel='stylesheet'>,
        the fetcher follows it AND parses Google Fonts URLs separately. Both
        paths must contribute to the extracted fonts.

        Pin the stylesheet response with a CORRECT Content-Type ("text/css")
        so the production content-type guard in _fetch_external_stylesheets
        actually allows the body through. Earlier version of this test
        passed only because Google Fonts URL parsing fired BEFORE urlopen
        and the external CSS body was silently rejected.
        """
        page_html = b"""
<!DOCTYPE html><html><head>
<link rel="stylesheet" href="/assets/style.css">
<link rel="stylesheet" href="https://fonts.googleapis.com/css2?family=Inter:wght@400">
</head><body><h1>Test</h1><p>Some content here.</p></body></html>
"""
        # NB: 'Helvetica Neue' is the marker that the EXTERNAL CSS path fired
        # (Inter would be found just from the Google Fonts query string).
        external_css = b"""
body { font-family: 'Helvetica Neue', sans-serif; color: #333333; }
h1 { font-family: 'Inter', sans-serif; color: #5A597A; }
"""

        def fake_urlopen(req, timeout=None):
            full_url = req.full_url if hasattr(req, "full_url") else str(req)
            if "style.css" in full_url:
                return _StubResponse(external_css, content_type="text/css")
            return _StubResponse(page_html, content_type="text/html")

        with mock.patch("extract.web_pdf._network_open", fake_urlopen):
            draft = extract_from_url("https://example.com/")

        # Fonts should be populated
        assert draft["fonts"], f"expected fonts to be populated, got {draft['fonts']}"
        font_names = {data.get("name", "") for data in draft["fonts"].values()}

        # BOTH paths must have fired:
        #   - Inter proves the Google Fonts query path fired (or the external CSS path did)
        #   - Helvetica Neue proves the external CSS fetch path specifically fired
        #     (it appears nowhere except in external_css)
        assert any("Helvetica" in n for n in font_names), \
            f"external CSS fetch path did not fire — expected 'Helvetica Neue' in fonts, got {font_names}"
        assert any("Inter" in n for n in font_names), \
            f"Google Fonts path did not fire — expected 'Inter' in fonts, got {font_names}"

    def test_synthetic_voice_artifact_validates(self):
        """Test 2: A synthetic voice extraction artifact (mirroring what real
        single-letter extraction would produce) should pass VoiceValidator
        (corpus size + section completeness). Validates that the Step 1.5
        prompt produces output the downstream consumers will accept."""
        from validator import VoiceValidator

        artifact_path = Path(__file__).parent / "fixtures" / "voice_extraction_synthetic_2026-01.yaml"
        assert artifact_path.exists(), "voice extraction artifact missing"

        voice = yaml.safe_load(artifact_path.read_text())

        # The artifact should be a valid populated voice section
        result = VoiceValidator.validate_voice(voice)
        assert result["status"] == "pass", \
            f"Synthetic voice artifact failed VoiceValidator: {result}"

        # And the corpus size check should specifically pass (>= 2000 words)
        assert result["checks"]["corpus"]["status"] == "pass"
        assert result["checks"]["corpus"]["word_count"] >= 2000

        # Section completeness should pass (positioning, tone, ≥2 core rules,
        # ≥3 anti-filler, ≥3 primary attributes)
        assert result["checks"]["completeness"]["status"] == "pass"

        # Spot-check that core_rules contain specific thesis vocabulary
        # (vs generic asset-management boilerplate):
        assert any("thesis" in r.lower() for r in voice["core_rules"]), \
            "core_rules should reference specific thesis vocabulary"

        # anti_filler should include institutional-finance buzzwords
        anti_filler_text = " ".join(voice["anti_filler"]).lower()
        assert "best-in-class" in anti_filler_text or "world-class" in anti_filler_text
        assert "leverage" in anti_filler_text

    def test_urllib_fallback_handles_total_network_failure(self, monkeypatch):
        """If everything fails, returns graceful empty draft (no exceptions)."""
        from subprocess import CompletedProcess

        def fake_run(*args, **kwargs):
            return CompletedProcess(args, returncode=1, stdout="")

        def fake_urlopen_fail(req, timeout=None):
            raise OSError("Network unreachable")

        with mock.patch("subprocess.run", fake_run), \
             mock.patch("extract.web_pdf._network_open", fake_urlopen_fail):
            draft = extract_from_url("https://0.0.0.0/")

        # Graceful degradation: empty colors/logos/fonts but valid structure
        assert "voice_corpus" in draft
        assert draft["voice_corpus"]["word_count"] == 0
        assert draft["source"]["type"] == "url"

    @pytest.mark.parametrize("url", [
        "http://127.0.0.1/private",
        "http://10.20.30.40/private",
        "http://169.254.169.254/latest/meta-data/",
        "http://[::1]/private",
        "http://[fe80::1]/private",
        "http://2130706433/private",
        "http://0x7f000001/private",
        "http://017700000001/private",
        "http://[::ffff:127.0.0.1]/private",
        "http://user@127.0.0.1/private",
        "http://example.com@127.0.0.1/private",
    ])
    def test_private_and_ambiguous_destinations_rejected_before_open(self, url):
        secret_url = f"{url}?token=do-not-log"
        with mock.patch("extract.web_pdf._network_open") as opener:
            draft = extract_from_url(secret_url)

        opener.assert_not_called()
        assert draft["error"] == "URL destination is not public"
        assert "do-not-log" not in str(draft)
        assert "?" not in draft["source"]["reference"]

    def test_hostname_with_any_private_dns_answer_is_rejected(self, monkeypatch):
        def mixed_getaddrinfo(host, port, *args, **kwargs):
            return [
                (2, 1, 6, "", ("93.184.216.34", port)),
                (2, 1, 6, "", ("10.0.0.8", port)),
            ]

        monkeypatch.setattr("socket.getaddrinfo", mixed_getaddrinfo)
        with mock.patch("extract.web_pdf._network_open") as opener:
            draft = extract_from_url("https://rebinding.example/brand")

        opener.assert_not_called()
        assert draft["error"] == "URL destination is not public"

    def test_redirect_location_is_validated_before_follow(self):
        from urllib.request import Request

        request = Request("https://example.com/start")

        def fake_build_opener(*handlers):
            # A real opener calls this hook before creating the redirected
            # request. The private Location must abort here, before open().
            handler = next(h for h in handlers if hasattr(h, "redirect_request"))
            handler.redirect_request(
                request, None, 302, "Found", {}, "http://127.0.0.1/admin"
            )
            raise AssertionError("private redirect was followed")

        with mock.patch("urllib.request.build_opener", fake_build_opener), \
             pytest.raises(ValueError, match="not public"):
            web_pdf_module._network_open(request, timeout=1)

    def test_https_connection_dials_pinned_ip_but_verifies_url_hostname(self):
        import http.client

        context = mock.Mock()
        raw_socket = mock.Mock()
        context.wrap_socket.return_value = mock.Mock()
        with mock.patch("socket.create_connection", return_value=raw_socket) as dial:
            connection = web_pdf_module._make_pinned_connection(
                http.client.HTTPSConnection,
                "brand.example",
                "93.184.216.34",
                context=context,
                timeout=3,
            )
            connection.connect()

        dial.assert_called_once_with(("93.184.216.34", 443), 3, None)
        context.wrap_socket.assert_called_once_with(
            raw_socket, server_hostname="brand.example"
        )

    def test_redirect_to_private_destination_is_rejected_before_body_read(self):
        response = _StubResponse(
            b"private response must never be consumed",
            final_url="http://127.0.0.1/admin?token=redirect-secret",
        )
        with mock.patch("extract.web_pdf._network_open", return_value=response):
            draft = extract_from_url("https://example.com/start?api_key=request-secret")

        assert response.read_calls == 0
        assert not draft["colors"]
        assert draft["error"] == "URL destination is not public"
        assert "request-secret" not in str(draft)
        assert "redirect-secret" not in str(draft)

    def test_public_mocked_destination_still_extracts_without_query_secret(self):
        response = _StubResponse(_BRAND_STUB_HTML, final_url="https://example.com/final")
        with mock.patch("extract.web_pdf._network_open", return_value=response):
            draft = extract_from_url("https://example.com/brand?token=do-not-persist")

        assert response.read_calls >= 1
        assert draft["colors"]["primary"]["hex"] == "#5A597A"
        assert "do-not-persist" not in str(draft)
