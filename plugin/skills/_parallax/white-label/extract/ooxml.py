"""OOXML theme XML parsing + PPTX/DOCX extraction.

PPTX and DOCX share the same theme XML structure (drawingml namespace), so
the parser is shared. Theme XML provides canonical brand declarations
(colors, fonts) at high confidence (0.9). Body text is aggregated for the
voice corpus; voice extraction itself is LLM-driven in SKILL.md.
"""

from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List

from .voice import _voice_corpus_from_text


_OOXML_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"



def _parse_pptx_master_typography(zf) -> tuple[Dict[str, Dict[str, Any]], Dict[str, float]]:
    from xml.etree import ElementTree as ET
    try:
        with zf.open("ppt/slideMasters/slideMaster1.xml") as f:
            root = ET.fromstring(f.read())
    except KeyError:
        return {}, {}
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main", "a": _OOXML_A_NS}
    txStyles = root.find(".//p:txStyles", ns)
    if txStyles is None: return {}, {}
    
    typography = {}
    confidences = {}
    
    # PPTX bodyStyle levels lvl1..lvl6 are *indent depths* inside the body
    # placeholder, not a heading hierarchy. lvl1 maps reasonably to "h2" (top
    # of body), but lvl3..lvl5 are progressively-indented bullet styles that
    # rarely correspond to a designer's intended h4/h5 typography. Cap their
    # confidence at 0.6 so downstream consumers can soft-weight or omit them.
    LOW_CONFIDENCE_LEVELS = {"h4", "h5", "body-md", "body-sm"}

    def parse_pr(pr_elem, level_name):
        if pr_elem is None: return
        defRPr = pr_elem.find(".//a:defRPr", ns)
        if defRPr is None: return
        sz = defRPr.get("sz")
        b = defRPr.get("b")
        lnSpc = pr_elem.find(".//a:lnSpc/a:spcPct", ns)

        # DESIGN.md spec only permits px / rem / em units; pt is rejected by the
        # linter. Convert pt → px at 96dpi (1pt = 4/3 px) and round to integer.
        # letterSpacing must be a dimension (with unit), not bare "0".
        style = {"letterSpacing": "0em"}
        conf = 0.5
        if sz:
            pt = int(sz) / 100
            px = round(pt * 4 / 3)
            style["fontSize"] = f"{px}px"
            # Top-of-hierarchy levels (titleStyle → h1, bodyStyle lvl1 → h2,
            # lvl2 → h3) are confident size signals. Deeper indent levels get
            # a lower ceiling — they're bullet styles, not heading scale.
            conf = 0.6 if level_name in LOW_CONFIDENCE_LEVELS else 0.85
        # DrawingML accepts xsd:boolean: "1"/"true" → bold, "0"/"false"/missing
        # → not bold. Symmetric with the DOCX path (line ~125).
        b_truthy = b is not None and str(b).lower() in ("1", "true")
        style["fontWeight"] = 700 if b_truthy else 400
        if lnSpc is not None and lnSpc.get("val"):
            style["lineHeight"] = f"{int(lnSpc.get('val'))/100000:g}"
            if conf < 0.7: conf = 0.7

        typography[level_name] = style
        confidences[f"typography.{level_name}"] = conf

    titleStyle = txStyles.find("p:titleStyle", ns)
    if titleStyle is not None:
        parse_pr(titleStyle.find("a:lvl1pPr", ns), "h1")
    
    bodyStyle = txStyles.find("p:bodyStyle", ns)
    if bodyStyle is not None:
        parse_pr(bodyStyle.find("a:lvl1pPr", ns), "h2")
        parse_pr(bodyStyle.find("a:lvl2pPr", ns), "h3")
        parse_pr(bodyStyle.find("a:lvl3pPr", ns), "h4")
        parse_pr(bodyStyle.find("a:lvl4pPr", ns), "h5")
        parse_pr(bodyStyle.find("a:lvl5pPr", ns), "body-md")
        parse_pr(bodyStyle.find("a:lvl6pPr", ns), "body-sm")
        
    return typography, confidences

def _parse_docx_style_typography(zf) -> tuple[Dict[str, Dict[str, Any]], Dict[str, float]]:
    from xml.etree import ElementTree as ET
    try:
        with zf.open("word/styles.xml") as f:
            root = ET.fromstring(f.read())
    except KeyError:
        return {}, {}
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    typography = {}
    confidences = {}
    
    name_map = {
        "heading 1": "h1", "heading 2": "h2", "heading 3": "h3",
        "heading 4": "h4", "heading 5": "h5", "normal": "body-md",
        "no spacing": "body-sm"
    }
    for style in root.findall("w:style", ns):
        name_elem = style.find("w:name", ns)
        if name_elem is None: continue
        val = name_elem.get(f"{{{ns['w']}}}val")
        if not val: continue
        name_val = val.lower()
        if name_val not in name_map: continue
        level_name = name_map[name_val]
        
        rPr = style.find("w:rPr", ns)
        sz = rPr.find("w:sz", ns) if rPr is not None else None
        b = rPr.find("w:b", ns) if rPr is not None else None
        
        pPr = style.find("w:pPr", ns)
        lnSpc = None
        if pPr is not None:
            spacing = pPr.find("w:spacing", ns)
            if spacing is not None:
                lnSpc = spacing.get(f"{{{ns['w']}}}line")
                
        # DESIGN.md spec only permits px / rem / em units; pt is rejected by the
        # linter. DOCX w:sz is half-points → pt → px at 4/3 ratio (96dpi).
        # letterSpacing must carry a unit.
        style_dict = {"letterSpacing": "0em"}
        conf = 0.5
        if sz is not None:
            val_sz = sz.get(f"{{{ns['w']}}}val")
            if val_sz:
                pt = int(val_sz) / 2
                px = round(pt * 4 / 3)
                style_dict["fontSize"] = f"{px}px"
                conf = 0.85
        # OOXML allows explicit <w:b w:val="0"/> or <w:b w:val="false"/> to
        # DISABLE bold inheritance. The element's presence alone isn't enough —
        # treat val="0" / "false" as 400, otherwise 700.
        if b is not None:
            b_val = (b.get(f"{{{ns['w']}}}val") or "").lower()
            style_dict["fontWeight"] = 400 if b_val in ("0", "false") else 700
        else:
            style_dict["fontWeight"] = 400
        if lnSpc:
            style_dict["lineHeight"] = f"{int(lnSpc)/240:g}"
            if conf < 0.7: conf = 0.7
            
        typography[level_name] = style_dict
        confidences[f"typography.{level_name}"] = conf
    return typography, confidences

def _detect_corner_radii_pptx(zf) -> Dict[str, str]:
    from xml.etree import ElementTree as ET
    import re
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main", "a": _OOXML_A_NS}
    radii = []
    for name in zf.namelist():
        if re.match(r"^ppt/slides/slide\d+\.xml$", name):
            with zf.open(name) as f:
                try:
                    root = ET.fromstring(f.read())
                except Exception:
                    continue
                for geom in root.findall(".//a:prstGeom[@prst='roundRect']", ns):
                    adj = geom.find("a:avLst/a:gd[@name='adj']", ns)
                    if adj is not None and adj.get("fmla"):
                        fmla = adj.get("fmla")
                        if fmla.startswith("val "):
                            try:
                                # OOXML `adj` for prstGeom roundRect is a
                                # PERCENTAGE in units of 1/100000 — value 50000
                                # means radius equals 50% of half the shorter
                                # side (i.e. full pill). It is NOT a pixel value
                                # and dividing by 1000 was nonsense. Without the
                                # parent shape's xfrm cx/cy we can't compute the
                                # absolute px radius, so we keep the raw
                                # percentage and bucket later.
                                radii.append(int(fmla[4:]))
                            except ValueError:
                                pass
    if not radii:
        return {}
    # Bucket by percentage tier rather than guessed px. Emit a fixed
    # representative px per tier so downstream consumers get a usable token.
    # adj < 10000 (<10%)              → sm  (4px)
    # adj 10000–25000 (10–25%)        → md  (8px)
    # adj 25000–45000 (25–45%)        → lg  (16px)
    # adj >= 45000 (≥45%, near-pill)  → full (9999px)
    counts: Dict[str, int] = {}
    for r in radii:
        if r < 10000:
            key = "sm"
        elif r < 25000:
            key = "md"
        elif r < 45000:
            key = "lg"
        else:
            key = "full"
        counts[key] = counts.get(key, 0) + 1

    tier_to_px = {"sm": "4px", "md": "8px", "lg": "16px", "full": "9999px"}
    return {key: tier_to_px[key] for key in counts}


def _parse_ooxml_theme(theme_xml_bytes: bytes) -> Dict[str, Any]:
    """Parse OOXML theme XML and return color scheme + font scheme.

    Returns colors keyed by OOXML slot name (dk1, lt1, dk2, lt2, accent1..accent6,
    hlink, folHlink) and fonts keyed by role (major, minor).
    """
    from xml.etree import ElementTree as ET

    ns = {"a": _OOXML_A_NS}
    root = ET.fromstring(theme_xml_bytes)

    colors: Dict[str, str] = {}
    clr_scheme = root.find(".//a:clrScheme", ns)
    if clr_scheme is not None:
        for slot in clr_scheme:
            slot_name = slot.tag.split("}", 1)[-1]
            srgb = slot.find("a:srgbClr", ns)
            sys = slot.find("a:sysClr", ns)
            if srgb is not None and srgb.get("val"):
                colors[slot_name] = "#" + srgb.get("val").upper()
            elif sys is not None and sys.get("lastClr"):
                colors[slot_name] = "#" + sys.get("lastClr").upper()

    fonts: Dict[str, str] = {}
    font_scheme = root.find(".//a:fontScheme", ns)
    if font_scheme is not None:
        for role_elem_name, role_key in [("majorFont", "major"), ("minorFont", "minor")]:
            role_elem = font_scheme.find(f"a:{role_elem_name}", ns)
            if role_elem is not None:
                latin = role_elem.find("a:latin", ns)
                if latin is not None and latin.get("typeface"):
                    fonts[role_key] = latin.get("typeface")

    return {"colors": colors, "fonts": fonts}


def _theme_to_role_map(theme: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    """Map OOXML theme slots to white-label color roles.

    accent1 → primary, accent2 → secondary, accent3 → accent.
    dk1 → text, lt1 → background. Confidence 0.9 (theme XML is canonical).
    """
    raw = theme.get("colors", {})
    role_colors: Dict[str, Dict[str, Any]] = {}

    slot_to_role = [
        ("accent1", "primary"),
        ("accent2", "secondary"),
        ("accent3", "accent"),
        ("dk1", "text"),
        ("lt1", "background"),
    ]
    for slot, role in slot_to_role:
        if slot in raw:
            role_colors[role] = {"hex": raw[slot], "confidence": 0.9}

    fonts: Dict[str, Dict[str, Any]] = {}
    raw_fonts = theme.get("fonts", {})
    if "major" in raw_fonts:
        fonts["header"] = {"name": raw_fonts["major"], "confidence": 0.9}
    if "minor" in raw_fonts:
        fonts["body"] = {"name": raw_fonts["minor"], "confidence": 0.9}

    return {"colors": role_colors, "fonts": fonts}


def extract_from_pptx(pptx_path: str) -> Dict[str, Any]:
    """Extract branding (visual + voice corpus) from a PowerPoint file.

    Visual: theme XML provides precise colors/fonts (high confidence).
    Voice corpus: aggregated slide body text for LLM voice extraction in SKILL.md.
    """
    pptx_file = Path(pptx_path)
    base_return = {
        "colors": {},
        "logos": {},
        "fonts": {},
        "source": {"type": "pptx", "reference": pptx_path},
        "extracted_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "confidence_scores": {},
        "voice_corpus": {"text": "", "word_count": 0, "truncated": False},
    }

    if not pptx_file.exists():
        return {**base_return, "error": "PPTX file not found"}

    try:
        import zipfile
        theme: Dict[str, Any] = {"colors": {}, "fonts": {}}
        typography: Dict[str, Any] = {}
        typo_conf: Dict[str, float] = {}
        rounded: Dict[str, str] = {}
        # Single open — theme parse, typography parse, and radii detection all
        # operate on the same archive (previously opened twice).
        with zipfile.ZipFile(pptx_path, "r") as zf:
            theme_paths = [n for n in zf.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")]
            if theme_paths:
                with zf.open(theme_paths[0]) as f:
                    theme = _parse_ooxml_theme(f.read())
            typography, typo_conf = _parse_pptx_master_typography(zf)
            rounded = _detect_corner_radii_pptx(zf)

        role_map = _theme_to_role_map(theme)

        body_text = ""
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            chunks: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(r.text for r in para.runs).strip()
                            if line:
                                chunks.append(line)
            body_text = "\n".join(chunks)
        except Exception:
            pass

        corpus = _voice_corpus_from_text(body_text)

        confidence_scores: Dict[str, float] = {}
        for role, data in role_map["colors"].items():
            confidence_scores[f"color_{role}"] = data["confidence"]
        for usage, data in role_map["fonts"].items():
            confidence_scores[f"font_{usage}"] = data["confidence"]

        if typography:
            confidence_scores.update(typo_conf)
        if rounded:
            confidence_scores["rounded"] = 0.5
            
        ret = {
            **base_return,
            "colors": role_map["colors"],
            "fonts": role_map["fonts"],
            "confidence_scores": confidence_scores,
            "voice_corpus": corpus,
        }
        if typography:
            ret["typography"] = typography
        if rounded:
            ret["rounded"] = rounded
        return ret

    except Exception as e:
        return {**base_return, "error": str(e)}


def extract_from_docx(docx_path: str) -> Dict[str, Any]:
    """Extract branding (visual + voice corpus) from a Word document.

    Visual: theme XML if present (high confidence), styles fallback otherwise.
    Voice corpus: aggregated paragraph text for LLM voice extraction in SKILL.md.
    """
    docx_file = Path(docx_path)
    base_return = {
        "colors": {},
        "logos": {},
        "fonts": {},
        "source": {"type": "docx", "reference": docx_path},
        "extracted_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "confidence_scores": {},
        "voice_corpus": {"text": "", "word_count": 0, "truncated": False},
    }

    if not docx_file.exists():
        return {**base_return, "error": "DOCX file not found"}

    try:
        import zipfile
        theme: Dict[str, Any] = {"colors": {}, "fonts": {}}
        typography: Dict[str, Any] = {}
        typo_conf: Dict[str, float] = {}
        rounded: Dict[str, str] = {}
        # Single open — theme parse and typography parse share the archive.
        with zipfile.ZipFile(docx_path, "r") as zf:
            theme_paths = [n for n in zf.namelist() if n.startswith("word/theme/") and n.endswith(".xml")]
            if theme_paths:
                with zf.open(theme_paths[0]) as f:
                    theme = _parse_ooxml_theme(f.read())
            typography, typo_conf = _parse_docx_style_typography(zf)

        role_map = _theme_to_role_map(theme)


        body_text = ""
        try:
            from docx import Document
            doc = Document(docx_path)
            chunks: List[str] = []
            for para in doc.paragraphs:
                line = para.text.strip()
                if line:
                    chunks.append(line)
            body_text = "\n".join(chunks)
        except Exception:
            pass

        corpus = _voice_corpus_from_text(body_text)

        confidence_scores: Dict[str, float] = {}
        for role, data in role_map["colors"].items():
            confidence_scores[f"color_{role}"] = data["confidence"]
        for usage, data in role_map["fonts"].items():
            confidence_scores[f"font_{usage}"] = data["confidence"]

        if typography:
            confidence_scores.update(typo_conf)
        
            
        ret = {
            **base_return,
            "colors": role_map["colors"],
            "fonts": role_map["fonts"],
            "confidence_scores": confidence_scores,
            "voice_corpus": corpus,
        }
        if typography:
            ret["typography"] = typography
        
        return ret

    except Exception as e:
        return {**base_return, "error": str(e)}
